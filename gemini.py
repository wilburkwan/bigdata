import os
import gradio as gr
import google.generativeai as genai
import PyPDF2
import docx
import pptx
import pandas as pd
import json

# 設定 API Key 環境變量
api_key = "AIzaSyBCjcuoP7Oh5eK5Dh2nk9urYVeFjsPcLjw"
os.environ['GIMINI_API_KEY'] = api_key

# 定義函數讀取不同格式的文件
def read_file(file_path):
    ext = os.path.splitext(file_path)[-1].lower()
    if ext == '.pdf':
        return read_pdf(file_path)
    elif ext == '.txt':
        return read_txt(file_path)
    elif ext == '.docx':
        return read_docx(file_path)
    elif ext == '.pptx':
        return read_pptx(file_path)
    elif ext in ['.xls', '.xlsx']:
        return read_excel(file_path)
    else:
        return "Unsupported file format."

def read_pdf(file_path):
    with open(file_path, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        text = ''
        for page_num in range(len(reader.pages)):
            page = reader.pages[page_num]
            text += page.extract_text()
    return text

def read_txt(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        text = file.read()
    return text

def read_docx(file_path):
    doc = docx.Document(file_path)
    text = ''
    for para in doc.paragraphs:
        text += para.text + '\n'
    return text

def read_pptx(file_path):
    prs = pptx.Presentation(file_path)
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + '\n'
    return text

def read_excel(file_path):
    df = pd.read_excel(file_path)
    return df.to_string()

# 設定 generative ai 的 API Key
genai.configure(api_key=os.environ['GIMINI_API_KEY'])
model = genai.GenerativeModel('gemini-pro')

# 定義聊天機器人的處理函數
def chat_with_ai(chat_history, user_input, instruction, after_prompt):
    # 讀取 a.txt 文件內容
    file_text = read_txt('data\mco2_transfer_history.csv')
    # 設定分析的 prompt
    chat_history_text = '\n'.join([f'User: {msg[0]}\nAI: {msg[1]}' for msg in chat_history])
    final_prompt = instruction + file_text + "，這是你要閱讀的文件，以下是你需要回答別人的問題（盡量不要超過100字。）" + user_input + after_prompt + chat_history_text
    
    response = model.generate_content(final_prompt)  # Assuming this method call returns the required response
    
    if response and 'candidates' in response:
        if response['candidates']:
            generated_text = response['candidates'][0]['content']['parts'][0]['text']
            return generated_text
        else:
            return "No candidates found in response."
    else:
        # 假設 response 是一個 GenerateContentResponse 物件，這裡直接將其轉換為字串
        response_str = str(response)  # 將 response 轉換為字串

        # 提取 JSON 部分的字串（此步驟根據實際 response 的格式可能需要調整）
        start_index = response_str.find("{")
        end_index = response_str.rfind("}") + 1
        json_str = response_str[start_index:end_index]

        # 解析 JSON 資料
        data = json.loads(json_str)
        text = data["candidates"][0]["content"]["parts"][0]["text"]

        return text

# 定義 Gradio 的界面
with gr.Blocks() as demo:
    gr.Markdown("# 數據分析師機器人")
    chatbot = gr.Chatbot()
    message = gr.Textbox(label="輸入您的問題")
    send_button = gr.Button("發送")

    def respond(message, chat_history):
        instruction = "現在你是數據分析師"
        after_prompt = "你需要回答別人的問題（盡量不要超過100字。），請換行排版後輸出"
        bot_response = chat_with_ai(chat_history, message, instruction, after_prompt)
        chat_history.append((message, bot_response))
        return chat_history, ""

    send_button.click(respond, inputs=[message, chatbot], outputs=[chatbot, message])

# 啟動 Gradio 界面
demo.launch(share=True)
